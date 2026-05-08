from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colors ──────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0b, 0x1e, 0x45)
BLUE   = RGBColor(0x3B, 0x5A, 0xC6)
RED    = RGBColor(0xd4, 0x00, 0x2a)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF4, 0xF5, 0xF7)
GRAY   = RGBColor(0x55, 0x55, 0x55)
GREEN  = RGBColor(0x1D, 0x7A, 0x4A)
YELLOW = RGBColor(0xb4, 0x53, 0x09)
PURPLE = RGBColor(0x6d, 0x28, 0xd9)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Helper ───────────────────────────────────────────────────────────────
def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, radius=False):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, l, t, w, h)  # MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, space_before=0):
    from pptx.util import Pt
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — 표지
# ═══════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, NAVY)

# 왼쪽 빨간 포인트 바
add_rect(slide, Inches(0), Inches(0), Inches(0.18), Inches(7.5), RED)

# 메인 타이틀
t = slide.shapes.add_textbox(Inches(0.7), Inches(1.8), Inches(8), Inches(1.5))
tf = t.text_frame
p = tf.paragraphs[0]
run = p.add_run(); run.text = "🐢 SlowBro"
run.font.size = Pt(64); run.font.bold = True; run.font.color.rgb = WHITE

# 서브타이틀
t2 = slide.shapes.add_textbox(Inches(0.7), Inches(3.3), Inches(10), Inches(0.8))
tf2 = t2.text_frame
p2 = tf2.paragraphs[0]
run2 = p2.add_run()
run2.text = "어르신을 위한 접근성 티켓 예매 브라우저 — 기술 문서"
run2.font.size = Pt(22); run2.font.color.rgb = RGBColor(0xBB, 0xCC, 0xFF)

# 태그라인
t3 = slide.shapes.add_textbox(Inches(0.7), Inches(4.2), Inches(10), Inches(0.6))
tf3 = t3.text_frame
p3 = tf3.paragraphs[0]
run3 = p3.add_run()
run3.text = "\"천천히 해도 괜찮아요 🌿\""
run3.font.size = Pt(16); run3.font.color.rgb = RGBColor(0x88, 0xAA, 0xFF); run3.font.italic = True

# 우측 하단 메타
add_text(slide, "Electron 28  ·  Cheerio  ·  Web Speech API  ·  Vanilla JS",
         Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
         size=13, color=RGBColor(0x66, 0x77, 0x99))

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — 목차
# ═══════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, LIGHT)

add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), NAVY)
add_text(slide, "목차", Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
         size=28, bold=True, color=WHITE)

items = [
    ("01", "문제 정의", "어르신의 디지털 속도 격차와 SlowBro의 솔루션"),
    ("02", "Tech Stack", "사용 기술과 각 레이어의 역할"),
    ("03", "아키텍처", "Electron Main/Renderer 이중 프로세스 구조"),
    ("04", "시스템 플로우", "SlowBro 변환, 브라우저 모드, 예매 가드 흐름"),
    ("05", "ERD", "핵심 데이터 구조와 관계"),
    ("06", "IPC API 명세", "프로세스 간 통신 채널 전체 목록"),
]

colors = [BLUE, RED, GREEN, YELLOW, PURPLE, RGBColor(0x0b, 0x6e, 0x9e)]

for i, (num, title, desc) in enumerate(items):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.5 + row * 2.4)
    box = add_rect(slide, x, y, Inches(3.9), Inches(2.0), colors[i])
    add_text(slide, num, x + Inches(0.2), y + Inches(0.15), Inches(1), Inches(0.5),
             size=28, bold=True, color=RGBColor(0xFF,0xFF,0xFF))
    add_text(slide, title, x + Inches(0.2), y + Inches(0.65), Inches(3.5), Inches(0.5),
             size=16, bold=True, color=WHITE)
    add_text(slide, desc, x + Inches(0.2), y + Inches(1.1), Inches(3.5), Inches(0.7),
             size=11, color=RGBColor(0xCC, 0xDD, 0xFF), wrap=True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — 문제 정의
# ═══════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, LIGHT)

add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), RED)
add_text(slide, "01  문제 정의", Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
         size=26, bold=True, color=WHITE)

# 문제 카드들
problems = [
    ("⚡", "속도 압박", "30초 내 좌석 선택,\n팝업·광고·복잡한 레이아웃"),
    ("👀", "정보 과부하", "작은 글씨, 수십 개 버튼,\n어디를 눌러야 할지 모름"),
    ("❌", "예매 실패", "에러 메시지 이해 못하고\n시간 초과로 반복 실패"),
]

for i, (icon, title, desc) in enumerate(problems):
    x = Inches(0.5 + i * 4.1)
    add_rect(slide, x, Inches(1.4), Inches(3.8), Inches(2.4), NAVY)
    add_text(slide, icon, x + Inches(0.2), Inches(1.55), Inches(0.8), Inches(0.6), size=30)
    add_text(slide, title, x + Inches(0.2), Inches(2.1), Inches(3.4), Inches(0.5),
             size=16, bold=True, color=WHITE)
    add_text(slide, desc, x + Inches(0.2), Inches(2.6), Inches(3.4), Inches(0.9),
             size=12, color=RGBColor(0xBB, 0xCC, 0xFF), wrap=True)

# 화살표 → 솔루션
add_text(slide, "→ SlowBro 솔루션", Inches(0.5), Inches(4.1), Inches(12), Inches(0.6),
         size=20, bold=True, color=BLUE)

solutions = [
    "🐢  핵심 단계만 추려 1단계씩 안내 (Step-by-step guide)",
    "🎨  원본 사이트 디자인을 계승하며 불필요 요소 제거",
    "🔊  한국어 TTS로 각 단계 읽어주기",
    "🛡️  이중 클릭·도배 방지 Guard로 실수 예매 차단",
]
for i, s in enumerate(solutions):
    add_text(slide, s, Inches(0.5), Inches(4.75 + i * 0.52), Inches(12.3), Inches(0.45),
             size=14, color=RGBColor(0x22, 0x22, 0x22))

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Tech Stack
# ═══════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, LIGHT)

add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), BLUE)
add_text(slide, "02  Tech Stack", Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
         size=26, bold=True, color=WHITE)

stacks = [
    (NAVY,   "Electron 28",       "Main/Renderer 이중 프로세스\ncontextBridge 보안 IPC"),
    (GREEN,  "Cheerio",           "서버사이드 jQuery 파싱\n실제 사이트 HTML에서 데이터 추출"),
    (RED,    "Web Speech API",    "한국어 TTS (ko-KR)\n0.9× 속도로 단계별 안내"),
    (PURPLE, "BrowserView",       "Electron 내장 Chromium\n원본 사이트 그대로 임베드"),
    (YELLOW, "CSS Custom Props",  "--site-primary/accent/bg\n원본 사이트 테마 동적 적용"),
    (RGBColor(0x0b,0x6e,0x9e), "electron-builder", "Windows/Mac/Linux 패키징\nNSIS 인스톨러 생성"),
]

for i, (c, name, desc) in enumerate(stacks):
    col = i % 3; row = i // 3
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.4 + row * 2.5)
    add_rect(slide, x, y, Inches(4.0), Inches(2.2), c)
    add_text(slide, name, x+Inches(0.25), y+Inches(0.25), Inches(3.5), Inches(0.6),
             size=17, bold=True, color=WHITE)
    add_text(slide, desc, x+Inches(0.25), y+Inches(0.85), Inches(3.5), Inches(1.1),
             size=12, color=RGBColor(0xCC,0xDD,0xFF), wrap=True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — 아키텍처
# ═══════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, LIGHT)

add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), NAVY)
add_text(slide, "03  아키텍처", Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
         size=26, bold=True, color=WHITE)

# Main Process 박스
add_rect(slide, Inches(0.3), Inches(1.3), Inches(4.0), Inches(5.5), NAVY)
add_text(slide, "Main Process", Inches(0.3), Inches(1.35), Inches(4.0), Inches(0.5),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

main_items = ["index.js\n(Window, BrowserView)", "pageParser.js\n(Cheerio 파싱)", "antiAbuse.js\n(Guard)", "preload/index.js\n(contextBridge)"]
main_colors = [BLUE, GREEN, RED, PURPLE]
for i, (item, c) in enumerate(zip(main_items, main_colors)):
    add_rect(slide, Inches(0.5), Inches(1.9 + i*1.1), Inches(3.6), Inches(0.9), c)
    add_text(slide, item, Inches(0.6), Inches(1.95 + i*1.1), Inches(3.4), Inches(0.8),
             size=11, color=WHITE, wrap=True)

# IPC 화살표 영역
add_rect(slide, Inches(4.5), Inches(2.5), Inches(2.0), Inches(3.0), RGBColor(0xE0,0xE8,0xFF))
add_text(slide, "IPC\n(contextBridge)", Inches(4.5), Inches(2.7), Inches(2.0), Inches(1.0),
         size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_text(slide, "invoke\nsend\non", Inches(4.5), Inches(3.7), Inches(2.0), Inches(1.0),
         size=11, color=BLUE, align=PP_ALIGN.CENTER)

# Renderer Process 박스
add_rect(slide, Inches(6.7), Inches(1.3), Inches(4.0), Inches(5.5), RGBColor(0x1a, 0x3a, 0x6e))
add_text(slide, "Renderer Process", Inches(6.7), Inches(1.35), Inches(4.0), Inches(0.5),
         size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

renderer_items = ["index.html\n(UI 구조)", "app.js\n(상태·이벤트·TTS)", "main.css\n(반응형 스타일)", "Web Speech API\n(TTS 음성 안내)"]
renderer_colors = [BLUE, GREEN, YELLOW, PURPLE]
for i, (item, c) in enumerate(zip(renderer_items, renderer_colors)):
    add_rect(slide, Inches(6.9), Inches(1.9 + i*1.1), Inches(3.6), Inches(0.9), c)
    add_text(slide, item, Inches(7.0), Inches(1.95 + i*1.1), Inches(3.4), Inches(0.8),
             size=11, color=WHITE, wrap=True)

# BrowserView
add_rect(slide, Inches(11.0), Inches(1.3), Inches(2.0), Inches(2.5), RED)
add_text(slide, "BrowserView\n\n실제 사이트\n임베드", Inches(11.0), Inches(1.4), Inches(2.0), Inches(2.2),
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 파일 구조 설명
add_rect(slide, Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.35), RGBColor(0xDD,0xE4,0xFF))
add_text(slide, "src/main/  ·  src/preload/  ·  src/engine/  ·  src/guard/  ·  renderer/  ·  mock-sites/",
         Inches(0.3), Inches(6.97), Inches(12.7), Inches(0.3),
         size=11, color=NAVY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — 시스템 플로우
# ═══════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, LIGHT)

add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), GREEN)
add_text(slide, "04  시스템 플로우", Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
         size=26, bold=True, color=WHITE)

# Flow 1 — SlowBro 변환
add_rect(slide, Inches(0.3), Inches(1.2), Inches(4.0), Inches(5.8), NAVY)
add_text(slide, "🐢 SlowBro 변환 모드", Inches(0.3), Inches(1.25), Inches(4.0), Inches(0.45),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

flow1 = [
    ("1", "URL 입력"),
    ("2", "page:load (IPC)"),
    ("3", "Cheerio 파싱"),
    ("4", "Steps + Theme 반환"),
    ("5", "applyTheme()"),
    ("6", "renderStep() 순차 표시"),
]
for i, (n, label) in enumerate(flow1):
    add_rect(slide, Inches(0.55), Inches(1.78 + i*0.75), Inches(3.5), Inches(0.6), BLUE)
    add_text(slide, f"  {n}.  {label}", Inches(0.55), Inches(1.82 + i*0.75), Inches(3.5), Inches(0.55),
             size=12, color=WHITE)

# Flow 2 — 브라우저 모드
add_rect(slide, Inches(4.7), Inches(1.2), Inches(4.0), Inches(5.8), RGBColor(0x1a,0x3a,0x6e))
add_text(slide, "🌐 원본 브라우저 모드", Inches(4.7), Inches(1.25), Inches(4.0), Inches(0.45),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

flow2 = [
    ("1", "원본 보기 버튼"),
    ("2", "browser:open (IPC)"),
    ("3", "BrowserView 생성"),
    ("4", "setBounds 계산"),
    ("5", "URL 로드"),
    ("6", "네비게이션 이벤트 수신"),
]
for i, (n, label) in enumerate(flow2):
    add_rect(slide, Inches(4.95), Inches(1.78 + i*0.75), Inches(3.5), Inches(0.6), RGBColor(0x0b,0x6e,0x9e))
    add_text(slide, f"  {n}.  {label}", Inches(4.95), Inches(1.82 + i*0.75), Inches(3.5), Inches(0.55),
             size=12, color=WHITE)

# Flow 3 — 예매 Guard
add_rect(slide, Inches(9.1), Inches(1.2), Inches(4.0), Inches(5.8), RGBColor(0x3d,0x0c,0x12))
add_text(slide, "🛡️ 예매 가드", Inches(9.1), Inches(1.25), Inches(4.0), Inches(0.45),
         size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

flow3 = [
    ("1", "예매 버튼 클릭"),
    ("2", "checkClick (250ms 체크)"),
    ("3", "checkSession (중복 체크)"),
    ("✓", "정상 → 예매 진행"),
    ("✗", "이상 → 경고 메시지"),
    ("✗", "중복 → 재시도 차단"),
]
colors3 = [RED, RED, RED, GREEN, RGBColor(0xAA,0x00,0x00), RGBColor(0xAA,0x00,0x00)]
for i, ((n, label), c) in enumerate(zip(flow3, colors3)):
    add_rect(slide, Inches(9.35), Inches(1.78 + i*0.75), Inches(3.5), Inches(0.6), c)
    add_text(slide, f"  {n}.  {label}", Inches(9.35), Inches(1.82 + i*0.75), Inches(3.5), Inches(0.55),
             size=12, color=WHITE)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — ERD
# ═══════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, LIGHT)

add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), PURPLE)
add_text(slide, "05  ERD — 핵심 데이터 구조", Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
         size=26, bold=True, color=WHITE)

entities = [
    (NAVY,   "AppState",      ["currentUrl: string", "currentStep: number", "steps: Step[]", "theme: SiteTheme", "fontSize: number", "ttsEnabled: boolean", "mode: 'slow'|'orig'"]),
    (BLUE,   "Step",          ["id: string", "label: string", "type: 'select'|'confirm'|'done'", "options: string[]", "hint: string", "nextStep: string|null"]),
    (GREEN,  "SiteTheme",     ["primary: string", "accent: string", "bg: string", "name: string", "italic: boolean"]),
    (RED,    "GuardStore",    ["clickTimes: Map<session,number[]>", "sessionAttempts: Map<uid:eid,number>", "HUMAN_MIN_MS: 250", "MAX_ATTEMPTS: 1"]),
    (YELLOW, "ParsedResult",  ["steps: Step[]", "theme: SiteTheme|null", "title: string", "siteKey: string"]),
    (PURPLE, "BrowserBounds", ["x: number (0)", "y: number (navbar.bottom)", "width: number", "height: number"]),
]

for i, (c, name, fields) in enumerate(entities):
    col = i % 3; row = i // 3
    x = Inches(0.3 + col * 4.35)
    y = Inches(1.3 + row * 2.85)
    add_rect(slide, x, y, Inches(4.1), Inches(2.6), c)
    add_text(slide, name, x+Inches(0.15), y+Inches(0.1), Inches(3.8), Inches(0.45),
             size=15, bold=True, color=WHITE)
    # 구분선
    add_rect(slide, x+Inches(0.1), y+Inches(0.58), Inches(3.9), Inches(0.02), RGBColor(0xFF,0xFF,0xFF))
    for j, field in enumerate(fields):
        add_text(slide, f"  {field}", x+Inches(0.1), y+Inches(0.65 + j*0.3), Inches(3.9), Inches(0.28),
                 size=10, color=RGBColor(0xDD,0xEE,0xFF))

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — IPC API 명세 (1/2)
# ═══════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, LIGHT)

add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), RGBColor(0x0b,0x6e,0x9e))
add_text(slide, "06  IPC API 명세 (1/2) — 핵심 채널", Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
         size=26, bold=True, color=WHITE)

# 테이블 헤더
headers = ["채널", "방향", "파라미터", "반환값"]
col_w = [2.8, 1.5, 4.5, 4.0]
header_x = [0.3, 3.15, 4.7, 9.25]
for j, (h, x, w) in enumerate(zip(headers, header_x, col_w)):
    add_rect(slide, Inches(x), Inches(1.2), Inches(w), Inches(0.4), NAVY)
    add_text(slide, h, Inches(x+0.05), Inches(1.22), Inches(w-0.1), Inches(0.36),
             size=12, bold=True, color=WHITE)

apis = [
    ("page:load",         "invoke →",  "url: string",                    "{ steps, theme, title }"),
    ("guard:checkClick",  "invoke →",  "sessionId: string",              "{ ok: bool, reason? }"),
    ("guard:checkSession","invoke →",  "userId, eventId: string",        "{ ok: bool, reason? }"),
    ("tts:speak",         "send →",    "text: string",                   "—"),
    ("window:set-icon",   "send →",    "dataUrl: base64 PNG",            "—"),
    ("browser:open",      "invoke →",  "url: string",                    "{ ok: bool }"),
    ("browser:close",     "invoke →",  "—",                             "{ ok: bool }"),
    ("browser:navigate",  "invoke →",  "url: string",                    "{ ok: bool }"),
    ("browser:setBounds", "invoke →",  "{ x, y, width, height }",        "{ ok: bool }"),
]

row_colors = [BLUE, RED, RED, GREEN, PURPLE, RGBColor(0x0b,0x6e,0x9e),
              RGBColor(0x0b,0x6e,0x9e), RGBColor(0x0b,0x6e,0x9e), RGBColor(0x0b,0x6e,0x9e)]

for i, (ch, dir_, param, ret) in enumerate(apis):
    y = Inches(1.7 + i * 0.55)
    bg = RGBColor(0xF0,0xF4,0xFF) if i%2==0 else WHITE
    add_rect(slide, Inches(0.3), y, Inches(12.9), Inches(0.5), bg)
    for j, (val, x, w) in enumerate(zip([ch, dir_, param, ret], header_x, col_w)):
        col_c = row_colors[i] if j==0 else RGBColor(0x22,0x22,0x22)
        bold = j==0
        add_text(slide, val, Inches(x+0.07), y+Inches(0.05), Inches(w-0.1), Inches(0.42),
                 size=11, bold=bold, color=col_c if j==0 else RGBColor(0x22,0x22,0x22))

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — IPC 이벤트 + 마무리
# ═══════════════════════════════════════════════════════════════════════
slide = blank_slide(prs)
fill_bg(slide, LIGHT)

add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.1), RGBColor(0x0b,0x6e,0x9e))
add_text(slide, "06  IPC API 명세 (2/2) — Push 이벤트", Inches(0.5), Inches(0.2), Inches(12), Inches(0.7),
         size=26, bold=True, color=WHITE)

add_text(slide, "Main → Renderer 푸시 이벤트 (BrowserView 상태 전달)",
         Inches(0.5), Inches(1.2), Inches(12), Inches(0.45),
         size=14, bold=True, color=NAVY)

push_events = [
    ("browser:url-changed",   "event ←",  "url: string",                "현재 URL 업데이트"),
    ("browser:title-changed", "event ←",  "title: string",              "페이지 타이틀 업데이트"),
    ("browser:loading",       "event ←",  "isLoading: boolean",         "로딩 스피너 표시 여부"),
]

headers2 = ["채널", "방향", "페이로드", "설명"]
col_w2 = [3.0, 1.5, 4.0, 4.3]
header_x2 = [0.3, 3.35, 4.9, 9.0]
for j, (h, x, w) in enumerate(zip(headers2, header_x2, col_w2)):
    add_rect(slide, Inches(x), Inches(1.75), Inches(w), Inches(0.4), NAVY)
    add_text(slide, h, Inches(x+0.05), Inches(1.77), Inches(w-0.1), Inches(0.36),
             size=12, bold=True, color=WHITE)

for i, (ch, dir_, payload, desc) in enumerate(push_events):
    y = Inches(2.25 + i * 0.6)
    bg = RGBColor(0xF0,0xF4,0xFF) if i%2==0 else WHITE
    add_rect(slide, Inches(0.3), y, Inches(12.9), Inches(0.55), bg)
    for j, (val, x, w) in enumerate(zip([ch, dir_, payload, desc], header_x2, col_w2)):
        add_text(slide, val, Inches(x+0.07), y+Inches(0.07), Inches(w-0.1), Inches(0.42),
                 size=11, bold=(j==0), color=GREEN if j==0 else RGBColor(0x22,0x22,0x22))

# 마무리 카드
add_rect(slide, Inches(0.3), Inches(4.1), Inches(12.7), Inches(3.0), NAVY)
add_text(slide, "🐢 SlowBro — 핵심 가치 요약", Inches(0.5), Inches(4.25), Inches(12), Inches(0.5),
         size=18, bold=True, color=WHITE)

summary = [
    "• 어르신 접근성:  복잡한 티켓 사이트를 3~4 단계 가이드로 단순화",
    "• 브랜드 계승:  원본 사이트 색상(SITE_THEMES)으로 신뢰감 유지",
    "• 이중 모드:  SlowBro 변환 모드 + 원본 그대로 보기 (BrowserView)",
    "• 안전 예매:  250ms 클릭 가드 + 세션당 1회 예매 제한",
    "• 음성 안내:  Web Speech API ko-KR TTS로 각 단계 읽어주기",
]
for i, s in enumerate(summary):
    add_text(slide, s, Inches(0.5), Inches(4.85 + i * 0.42), Inches(12.3), Inches(0.38),
             size=13, color=RGBColor(0xBB, 0xCC, 0xFF))

# ── 저장 ─────────────────────────────────────────────────────────────
out_path = r"C:\khu\khuton\slowbro\demo\SlowBro_기술문서.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
